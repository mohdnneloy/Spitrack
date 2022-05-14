# =============== Imports ==================
import urllib
from bs4 import BeautifulSoup
import requests, PyPDF2
from io import BytesIO
import docx2txt
import urllib.request as urllib2
from pptx import Presentation


# =============== Method - Get URL List ==================

def getURLlist(root_url, strategies, depth):

    root_url = root_url
    strategies = strategies
    depth = depth

    # Receiving Lists
    pdf_urls = []
    word_urls = []
    powerpoint_urls = []
    text_urls = []
    non_html_text_urls = []

    # Main Extracted urls
    main_pdf_urls = []
    main_word_urls = []
    main_powerpoint_urls = []
    main_text_urls = []
    main_non_html_text_urls = []

    # Main Collected Urls
    all_urls = []

    # Collected Urls in a given depth for a given url in the main url list, consider it as temporary
    all_urls2 = []

    # Collecting urls from every depth
    for depthCount in range(depth):

        if depthCount == 0:
            all_urls.append({'url': root_url, 'depth': depthCount + 1})

        for murl in all_urls:
            print()
            print('Crawling URL Now: ', murl['url'], 'at depth: ', depthCount, '================================')

            try:
                response = requests.get(murl['url'], timeout=5)
                content = BeautifulSoup(response.text, 'lxml')
                all_urls_tags = content.find_all(['link', 'a'])

                # Selecting URLs which are secure "https" out of the crawled urls from a given URL in the main url list "murl"
                for url in all_urls_tags:
                    try:
                        if url['href'].startswith('https'):
                            all_urls2.append({'url': url['href'], 'depth': depthCount + 1})
                            print({'url': url['href'], 'depth': depthCount + 1})
                        # else:
                        #     print('--Url not Included:', url['href'])
                    except Exception as ex:
                        print('--Url not Included: ', url['href'], ex)

            except Exception as e:
                print('---Url not Crawled (Max tries Exceeded / Authentication Required): ', murl['url'])

        # Classifying collected urls from one depth in "all_urls"
        pdf_urls, word_urls, powerpoint_urls, text_urls, non_html_text_urls = (
            gerURLlistClassified(all_urls2, strategies))

        # Appending the new lists in the declared lists
        main_pdf_urls.extend(pdf_urls)
        main_word_urls.extend(word_urls)
        main_powerpoint_urls.extend(powerpoint_urls)
        main_text_urls.extend(text_urls)
        main_non_html_text_urls.extend(non_html_text_urls)

        # Passing all collected urls for the next depth crawl
        all_urls = all_urls2.copy()

    print()
    print('For ', root_url, " we got this specific URLS")
    print()
    print('PDF URLS ===========')
    for i in range(len(pdf_urls)):
        print(main_pdf_urls[i])
    print()

    print('Word URLS ===========')
    for i in range(len(word_urls)):
        print(word_urls[i])
    print()

    print('Powerpoint URLS ===========')
    for i in range(len(powerpoint_urls)):
        print(powerpoint_urls[i])
    print()

    print('Text URLS ===========')
    for i in range(len(text_urls)):
        print(text_urls[i])
    print()

    print('Non Html Text URLS ===========')
    for i in range(len(non_html_text_urls)):
        print(non_html_text_urls[i])
    print()

    return main_pdf_urls, main_word_urls, main_powerpoint_urls, main_text_urls, main_non_html_text_urls


# =============== Method - Get URL List Classified ==================

def gerURLlistClassified(all_urls, strategies):

    # Required Lists

    all_urls = all_urls
    strategies = strategies
    pdf_urls = []
    word_urls = []
    powerpoint_urls = []
    text_urls = []
    non_html_text_urls = []

    # Extracting urls with extension '.pdf', '.docx', '.pptx', '.txt', '.html', '.htm' from give list of urls

    for url in all_urls:
        try:
            if strategies['pdf'] & (url['url'].endswith('.pdf')):
                if not (url['url'] in pdf_urls):
                    pdf_urls.append({'url': url['url'], 'depth': url['depth']})

            if strategies['word'] & (url['url'].endswith('.docx')):
                if not (url['url'] in word_urls):
                    word_urls.append({'url': url['url'], 'depth': url['depth']})

            if strategies['powerpoint'] & (url['url'].endswith('.pptx')):
                if not (url['url'] in powerpoint_urls):
                    powerpoint_urls.append({'url': url['url'], 'depth': url['depth']})

            if strategies['text'] & (url['url'].endswith('.txt')):
                if not (url['url'] in text_urls):
                    text_urls.append({'url': url['url'], 'depth': url['depth']})

            if strategies['non_html_text'] & (url['url'].endswith('.html') | (url['url'].endswith('.htm'))):
                if not (url['url'] in text_urls):
                    non_html_text_urls.append({'url': url['url'], 'depth': url['depth']})

        except Exception as e:
            print('Error: ', e, 'for URL: ', url)
            # print('')

    # Returning all extracted url lists
    return pdf_urls, word_urls, powerpoint_urls, text_urls, non_html_text_urls


# =============== Method - CrawlAll-Data Based on Strategies ==================

def parseAll(pdf_urls, word_urls, powerpoint_urls, text_urls, non_html_text_urls):

    pdf_content = []
    word_content = []
    powerpoint_content = []
    text_content = []
    non_html_text_content = []

    # If PDF URL List not empty then parse data
    if pdf_urls:
        pdf_content = parsePDF(pdf_urls)

    # If Word URL List not empty then parse data
    if word_urls:
        word_content = parseWord(word_urls)

    # If Powerpoint URL List not empty then parse data
    if powerpoint_urls:
        powerpoint_content = parsePowerpoint(powerpoint_urls)

    # If Text URL List not empty then parse data
    if text_urls:
        text_content = parseText(text_urls)

    # If Non-Html-Text URL List not empty then parse data
    if non_html_text_urls:
        non_html_text_content = parseHtmlText(non_html_text_urls)

    return pdf_content, word_content, powerpoint_content, text_content, non_html_text_content


# =============== Method - Parse PDF ==================

def parsePDF(pdf_urls):

    pdf_content = []
    done = []
    for url in pdf_urls:
        try:
            if url not in done:
                per_page = []
                response = requests.get(url['url'], timeout=5)
                print('Entered')
                my_raw_data = response.content

                with BytesIO(my_raw_data) as data:
                    read_pdf = PyPDF2.PdfFileReader(data)

                    for page in range(read_pdf.getNumPages()):
                        # print(read_pdf.getPage(page).extractText())
                        text = read_pdf.getPage(page).extractText()
                        per_page.append(text)

                    string = ''
                    for text_object in per_page:
                        string = string + text_object + ' '
                print('Done PDF Parsing:', url['url'])
                pdf_content.append(string)
                done.append(url)
            # else:
            #     print("Already Parsed!")
        except Exception as e:
            print('skip')
            continue

    print('Done with PDF URLS')
    print()
    return pdf_content


# =============== Method - Parse Word ==================

def parseWord(word_urls):

    word_content = []
    done = []
    for url in word_urls:
        try:
            if url not in done:
                docx = BytesIO(requests.get(url['url'], timeout=5).content)
                print('Entered')
                text = docx2txt.process(docx)
                # print(text)
                print('Done Word Parsing:', url['url'])
                word_content.append(text)
                done.append(url)
            # else:
            #     print("Already Parsed!")
        except Exception as e:
            print('skip')
            continue

    print('Done with Word URLS')
    print()
    return word_content


# =============== Method - Parse Powerpoint ==================

def parsePowerpoint(powerpoint_urls):

    powerpoint_content = []
    done = []
    for url in powerpoint_urls:
        try:
            if url not in done:
                string = ''
                response = requests.get(url['url'], timeout=5)
                print('Entered')
                my_raw_data = response.content
                with BytesIO(my_raw_data) as data:
                    prs = Presentation(data)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                # print(shape.text)
                                string = string + shape.text + ' '
                print('Done Powerpoint Parsing:', url['url'])
                powerpoint_content.append(string)
                done.append(url)
            # else:
            #     print("Already Parsed!")
        except Exception as e:
            print('skip')
            continue

    print('Done with Powerpoint URLS')
    print()
    return powerpoint_content


# =============== Method - Parse Text ==================

def parseText(text_urls):

    text_content = []
    done =[]
    for url in text_urls:
        try:
            if url not in done:
                string = ''
                file = urllib.request.urlopen(url['url'], timeout=5)
                print('Entered')

                for line in file:
                    decoded_line = line.decode("utf-8")
                    # print(decoded_line)
                    string = string + decoded_line + ' '
                print('Done Text Parsing:', url['url'])
                text_content.append(string)
                done.append(url)
            # else:
            #     print("Already Parsed!")
        except Exception as e:
            print('skip error:', e)
            continue

    print('Done with Text URLS')
    print()
    return text_content


# =============== Method - Parse HTML Text ==================

def parseHtmlText(non_html_text_urls):

    non_html_content = []
    done = []
    for url in non_html_text_urls:
        try:
            if url not in done:
                string = ''
                html = urllib.request.urlopen(url['url'], timeout=5)
                print('Entered')
                htmlParse = BeautifulSoup(html, 'html.parser')
                for para in htmlParse.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'article']):
                    # print(para.get_text())
                    string = string + para.get_text() + ' '
                # print(string)
                print('Done HTML Parsing:', url['url'])
                non_html_content.append(string)
                done.append(url)
            # else:
            #     print("Already Parsed!")
        except Exception as e:
            print('skip')
            continue


    print('Done with Html URLS')
    print()
    return non_html_content